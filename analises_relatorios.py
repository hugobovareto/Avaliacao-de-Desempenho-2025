import os
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from docxtpl import DocxTemplate, InlineImage
from docx.shared import Mm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import warnings
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
warnings.filterwarnings('ignore')


# RELATÓRIOS INDIVIDUAIS EM PPT

# Carregando o dataframe processado
df_total = pd.read_csv('dados/processado_dados_avaliacao.csv')


# =================== CONFIGURAÇÕES ===================
DIRETORIO_BASE = r"D:\Scripts_Python\FGV\Avaliacoes_de_Desempenho_2025"
CAMINHO_TEMPLATE = os.path.join(DIRETORIO_BASE, "templates", "[GARN] Modelo de apresentação de slides.pptx")
DIRETORIO_SAIDA = os.path.join(DIRETORIO_BASE, "relatorios_gerados")
DIRETORIO_GRAFICOS_TEMP = os.path.join(DIRETORIO_BASE, "graficos_temp")

# Criar diretórios se não existirem
os.makedirs(DIRETORIO_SAIDA, exist_ok=True)
os.makedirs(DIRETORIO_GRAFICOS_TEMP, exist_ok=True)

# =================== DEFINIÇÕES DAS COLUNAS ===================
# Questões gerais
COLUNAS_GERAIS = [
    'Visao Sistemica_Identifica desafios',
    'Visao Sistematica_Compartilha ferramentas',
    'Visao Sistemica_Colabora com as demais frentes',
    'Visao Sistemica_Compreensao do impacto',
    'Gestao e Lideranca_Canal aberto',
    'Gestao e Lideranca_Incentiva e colabora',
    'Gestao e Lideranca_Autogestao de tempo',
    'Relacionamento_Cativa parceiros',
    'Relacionamento_Cria vinculos',
    'Comunicacao_Assertivo',
    'Comunicacao_Nao Violenta',
    'Comunicacao_Postura Profissional',
    'Aprendizagem_Compartilha experiencia',
    'Aprendizagem_Conhecimentos de fontes internas e externas',
    'Aprendizagem_Busca se desenvolver',
    'Aprendizagem_Autonomia',
    'Execucao_Resolve as dificuldades',
    'Execucao_Imprevistos e alteracoes',
    'Execucao_Ideias em Acoes'
]

# Questões de liderança
COLUNAS_LIDERANCA = [
    'Lideranca_Desenvolvimento de Pessoas',
    'Lideranca_Visão Estratégica',
    'Lideranca_Delegação',
    'Lideranca_Gerenciamento de Riscos',
    'Lideranca_Monitoramento de Resultados'
]

# Mapeamento de questões para placeholders gráficos
MAPEAMENTO_QUESTOES_GRAFICOS = {
    'Visao Sistemica_Identifica desafios': 'VS_1',
    'Visao Sistematica_Compartilha ferramentas': 'VS_2',
    'Visao Sistemica_Colabora com as demais frentes': 'VS_3',
    'Visao Sistemica_Compreensao do impacto': 'VS_4',
    'Gestao e Lideranca_Canal aberto': 'GL_1',
    'Gestao e Lideranca_Incentiva e colabora': 'GL_2',
    'Gestao e Lideranca_Autogestao de tempo': 'GL_3',
    'Relacionamento_Cativa parceiros': 'REL_1',
    'Relacionamento_Cria vinculos': 'REL_2',
    'Comunicacao_Assertivo': 'COM_1',
    'Comunicacao_Nao Violenta': 'COM_2',
    'Comunicacao_Postura Profissional': 'COM_3',
    'Aprendizagem_Compartilha experiencia': 'AD_1',
    'Aprendizagem_Conhecimentos de fontes internas e externas': 'AD_2',
    'Aprendizagem_Busca se desenvolver': 'AD_3',
    'Aprendizagem_Autonomia': 'AD_4',
    'Execucao_Resolve as dificuldades': 'EX_1',
    'Execucao_Imprevistos e alteracoes': 'EX_2',
    'Execucao_Ideias em Acoes': 'EX_3'
}

# Mapeamento de tipos para nomes amigáveis
TIPOS_AMIGAVEIS = {
    'autoavaliacao': 'Autoavaliação',
    'avaliacao_pelo_lider': 'Avaliação do Líder',
    'avaliacao_pelo_liderado': 'Avaliação dos Liderados'
}

CORES_TIPOS = {
    'autoavaliacao': '#00a6dc',  # Azul
    'avaliacao_pelo_lider': '#58a75b',  # Verde
    'avaliacao_pelo_liderado': '#cc8a42'  # Marrom
}

# =================== FUNÇÕES AUXILIARES ===================
def formatar_nota(valor):
    """Formata uma nota com 1 casa decimal"""
    if pd.isna(valor):
        return "Não se aplica"
    return f"{valor:.1f}".replace('.', ',')

def filtrar_dados_pessoa(df, nome):
    """Filtra o dataframe para uma pessoa específica"""
    return df[df['Nome'] == nome]

def calcular_medias_tipo(df_pessoa, tipo, colunas):
    """Calcula a média para um tipo específico e um conjunto de colunas"""
    df_tipo = df_pessoa[df_pessoa['Tipo'] == tipo]
    if df_tipo.empty:
        return None
    return df_tipo[colunas].mean()

def criar_grafico_barras_questao(df_pessoa, questao, tipo_placeholder):
    """Cria um gráfico de barras para uma questão específica"""
    # Coletar médias por tipo
    dados = {}
    for tipo in ['autoavaliacao', 'avaliacao_pelo_lider', 'avaliacao_pelo_liderado']:
        df_tipo = df_pessoa[df_pessoa['Tipo'] == tipo]
        if not df_tipo.empty and questao in df_tipo.columns:
            media = df_tipo[questao].mean()
            if not pd.isna(media):
                dados[tipo] = media
    
    if not dados:
        return None
    
    # Criar gráfico
    fig, ax = plt.subplots(figsize=(4, 3))
    tipos_presentes = list(dados.keys())
    valores = [dados[t] for t in tipos_presentes]
    cores = [CORES_TIPOS[t] for t in tipos_presentes]
    nomes = [TIPOS_AMIGAVEIS[t] for t in tipos_presentes]
    
    bars = ax.bar(range(len(tipos_presentes)), valores, color=cores)
    ax.set_ylim(0, 6)  # Escala de 0-6
    
    # Adicionar valores nas barras
    for bar, valor in zip(bars, valores):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{valor:.1f}', ha='center', va='bottom')
    
    ax.set_xticks(range(len(tipos_presentes)))
    ax.set_xticklabels(nomes, rotation=45, ha='right')
    ax.set_ylabel('Média')
    ax.set_title(f'{questao.split("_")[0]}', fontsize=10)
    
    plt.tight_layout()
    
    # Salvar gráfico
    caminho_grafico = os.path.join(DIRETORIO_GRAFICOS_TEMP, f"{tipo_placeholder}_{id(df_pessoa)}.png")
    fig.savefig(caminho_grafico, dpi=150, bbox_inches='tight')
    plt.close(fig)
    
    return caminho_grafico

def criar_radar_geral(df_pessoa):
    """Cria gráfico de radar para competências gerais"""
    # Definir categorias e questões
    categorias = {
        'Visão Sistêmica': [
            'Visao Sistemica_Identifica desafios',
            'Visao Sistematica_Compartilha ferramentas',
            'Visao Sistemica_Colabora com as demais frentes',
            'Visao Sistemica_Compreensao do impacto'
        ],
        'Gestão e Liderança': [
            'Gestao e Lideranca_Canal aberto',
            'Gestao e Lideranca_Incentiva e colabora',
            'Gestao e Lideranca_Autogestao de tempo'
        ],
        'Relacionamento': [
            'Relacionamento_Cativa parceiros',
            'Relacionamento_Cria vinculos'
        ],
        'Comunicação': [
            'Comunicacao_Assertivo',
            'Comunicacao_Nao Violenta',
            'Comunicacao_Postura Profissional'
        ],
        'Aprendizagem e Desenvolvimento': [
            'Aprendizagem_Compartilha experiencia',
            'Aprendizagem_Conhecimentos de fontes internas e externas',
            'Aprendizagem_Busca se desenvolver',
            'Aprendizagem_Autonomia'
        ],
        'Execução': [
            'Execucao_Resolve as dificuldades',
            'Execucao_Imprevistos e alteracoes',
            'Execucao_Ideias em Acoes'
        ]
    }
    
    # Calcular médias por categoria para cada tipo
    dados_radar = {}
    for tipo in ['autoavaliacao', 'avaliacao_pelo_lider', 'avaliacao_pelo_liderado']:
        df_tipo = df_pessoa[df_pessoa['Tipo'] == tipo]
        if not df_tipo.empty:
            medias_categoria = []
            for categoria, questoes in categorias.items():
                # Filtrar questões que existem no dataframe
                questoes_existentes = [q for q in questoes if q in df_tipo.columns]
                if questoes_existentes:
                    media = df_tipo[questoes_existentes].mean().mean()
                    medias_categoria.append(media if not pd.isna(media) else 0)
                else:
                    medias_categoria.append(0)
            dados_radar[tipo] = medias_categoria
    
    if not dados_radar:
        return None
    
    # Configurar radar
    N = len(categorias)
    angulos = [n / float(N) * 2 * np.pi for n in range(N)]
    angulos += angulos[:1]
    
    # Criar figura com mais espaço para a legenda
    fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(projection='polar'))
    
    # Plotar cada tipo
    for tipo in dados_radar.keys():
        valores = dados_radar[tipo]
        valores += valores[:1]
        ax.plot(angulos, valores, 'o-', linewidth=2, 
                label=TIPOS_AMIGAVEIS[tipo], color=CORES_TIPOS[tipo])
        ax.fill(angulos, valores, alpha=0.1, color=CORES_TIPOS[tipo])
    
    # Configurar
    ax.set_xticks(angulos[:-1])
    
    # Configurar labels com melhor espaçamento
    xtick_labels = list(categorias.keys())
    ax.set_xticklabels(xtick_labels, fontsize=11)
    
    # Ajustar posição das labels para ficarem fora do gráfico
    for label, angle in zip(ax.get_xticklabels(), angulos[:-1]):
        texto = label.get_text()

        # Regra geral baseada no ângulo
        if 0 <= angle < np.pi:
            alinhamento = 'left'
        else:
            alinhamento = 'right'

        # Exceções específicas
        if texto == 'Relacionamento':
            alinhamento = 'right'

        if texto == 'Execução':
            alinhamento = 'left'

        label.set_horizontalalignment(alinhamento)
        label.set_rotation(angle * 180 / np.pi - 90)


    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(['1', '2', '3', '4', '5'], fontsize=9)
    
    # Adicionar grade para melhor visualização
    ax.grid(True, alpha=0.3)
    
    # Posicionar legenda fora do gráfico
    ax.legend(loc='upper left', bbox_to_anchor=(-0.35, 1.25), 
              fontsize=10, framealpha=0.9)
    
    # Ajustar layout para dar mais espaço
    plt.tight_layout(rect=[0, 0, 1, 1]) 
    
    # Salvar gráfico
    caminho_grafico = os.path.join(DIRETORIO_GRAFICOS_TEMP, f"radar_geral_{id(df_pessoa)}.png")
    fig.savefig(caminho_grafico, dpi=150, bbox_inches='tight', pad_inches=0.5)
    plt.close(fig)
    
    return caminho_grafico

def criar_radar_lideranca(df_pessoa):
    """Cria gráfico de radar para competências de liderança"""
    # Calcular médias por questão para cada tipo
    dados_radar = {}
    for tipo in ['autoavaliacao', 'avaliacao_pelo_lider', 'avaliacao_pelo_liderado']:
        df_tipo = df_pessoa[df_pessoa['Tipo'] == tipo]
        if not df_tipo.empty:
            medias = []
            for questao in COLUNAS_LIDERANCA:
                if questao in df_tipo.columns:
                    media = df_tipo[questao].mean()
                    medias.append(media if not pd.isna(media) else 0)
                else:
                    medias.append(0)
            dados_radar[tipo] = medias
    
    if not dados_radar:
        return None
    
    # Configurar radar
    N = len(COLUNAS_LIDERANCA)
    angulos = [n / float(N) * 2 * np.pi for n in range(N)]
    angulos += angulos[:1]
    
    # Criar figura com mais espaço para a legenda
    fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(projection='polar'))
    
    # Encurtar nomes para o radar
    labels = [q.replace('Lideranca_', '').replace('_', ' ') for q in COLUNAS_LIDERANCA]
    
    # Plotar cada tipo
    for tipo in dados_radar.keys():
        valores = dados_radar[tipo]
        valores += valores[:1]
        ax.plot(angulos, valores, 'o-', linewidth=2, 
                label=TIPOS_AMIGAVEIS[tipo], color=CORES_TIPOS[tipo])
        ax.fill(angulos, valores, alpha=0.1, color=CORES_TIPOS[tipo])
    
    # Configurar
    ax.set_xticks(angulos[:-1])
    
    # Configurar labels com melhor espaçamento
    ax.set_xticklabels(labels, fontsize=10)
    
    # Ajustar posição das labels para ficarem fora do gráfico
    for label, angle in zip(ax.get_xticklabels(), angulos[:-1]):
        texto = label.get_text()
    
        # Regra geral baseada no ângulo
        if 0 <= angle < np.pi:
            alinhamento = 'left'
        else:
            alinhamento = 'right'

        # Exceções específicas
        if texto == 'Delegação':
            alinhamento = 'right'

        label.set_horizontalalignment(alinhamento)
        label.set_rotation(angle * 180 / np.pi - 90)


    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(['1', '2', '3', '4', '5'], fontsize=9)
    
    # Adicionar grade para melhor visualização
    ax.grid(True, alpha=0.3)
    
    # Posicionar legenda fora do gráfico
    ax.legend(loc='upper left', bbox_to_anchor=(-0.35, 1.25), 
              fontsize=10, framealpha=0.9)
    
    # Ajustar layout para dar mais espaço
    plt.tight_layout()  
    
    # Salvar gráfico
    caminho_grafico = os.path.join(DIRETORIO_GRAFICOS_TEMP, f"radar_lideranca_{id(df_pessoa)}.png")
    fig.savefig(caminho_grafico, dpi=150, bbox_inches='tight', pad_inches=0.5)
    plt.close(fig)
    
    return caminho_grafico


def substituir_respostas_abertas(slide, placeholder_tag, lista_respostas, texto_vazio="Não informado"):
    """
    Substitui um placeholder por uma lista de respostas abertas, cada uma em um parágrafo.
    """
    # Procurar pelo placeholder no slide
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.has_text_frame:
            continue
            
        if f"{{{{{placeholder_tag}}}}}" in shape.text:
            # Limpar o texto atual
            shape.text = ""
            
            # Adicionar o novo texto
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            # Se não houver respostas, usar texto_vazio
            if not lista_respostas or len(lista_respostas) == 0:
                p = text_frame.add_paragraph()
                p.text = texto_vazio
                p.alignment = PP_ALIGN.CENTER
                p.space_before = Pt(0)
                
                for run in p.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(0xC4, 0x80, 0x3F)  # #C4803F
                    run.font.name = 'Nunito'
            else:
                # Adicionar cada resposta como um parágrafo separado
                for i, resposta in enumerate(lista_respostas):
                    p = text_frame.add_paragraph()
                    p.text = str(resposta).strip()
                    p.alignment = PP_ALIGN.LEFT  # Respostas alinhadas à esquerda
                    p.space_before = Pt(0)
                    
                    # Formatar
                    for run in p.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Preto para respostas
                        run.font.name = 'Nunito'
                    
                    # Adicionar espaço entre parágrafos (exceto no último)
                    if i < len(lista_respostas) - 1:
                        p.space_after = Pt(12)  # Espaço maior entre respostas
            
            # Configurar margens
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            return True
    
    return False



def limpar_placeholders(slide):
    """Remove os placeholders {{...}} do slide"""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if "{{" in shape.text:
                shape.text = shape.text.replace("{{", "").replace("}}", "")

def adicionar_imagem_no_placeholder(slide, placeholder_nome, caminho_imagem, posicao=None, tamanho=None):
    """Adiciona uma imagem no lugar de um placeholder"""
    for shape in slide.shapes:
        if hasattr(shape, "text") and f"{{{{{placeholder_nome}}}}}" in shape.text:
            # Se encontramos um placeholder de texto
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            # Remover a forma de texto
            sp = shape._element
            sp.getparent().remove(sp)
            
            # Adicionar imagem
            slide.shapes.add_picture(caminho_imagem, left, top, width, height)
            return True
    return False

def substituir_texto_no_slide(slide, placeholder, valor, max_caracteres=None):
    """Substitui texto em um placeholder"""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if f"{{{{{placeholder}}}}}" in shape.text:
                # Limitar texto se max_caracteres for especificado
                if max_caracteres and len(str(valor)) > max_caracteres:
                    valor = str(valor)[:max_caracteres-3] + "..."
                shape.text = shape.text.replace(f"{{{{{placeholder}}}}}", str(valor))
                return True
    return False

def substituir_texto_formatado(slide, placeholder, valor, tamanho_fonte=14, cor_hex="#CC8A42", centralizar=True, negrito=True):
    """Substitui texto em um placeholder com formatação específica"""
    for shape in slide.shapes:
        if hasattr(shape, "text") and f"{{{{{placeholder}}}}}" in shape.text:
            shape.text = shape.text.replace(f"{{{{{placeholder}}}}}", str(valor))
            
            # Formatação do texto
            for paragraph in shape.text_frame.paragraphs:
                if centralizar:
                    paragraph.alignment = PP_ALIGN.CENTER
                
                for run in paragraph.runs:
                    run.font.size = Pt(tamanho_fonte)
                    run.font.bold = negrito
                    
                    # Converter cor hexadecimal para RGB
                    if cor_hex:
                        cor_hex = cor_hex.lstrip('#')
                        r = int(cor_hex[0:2], 16)
                        g = int(cor_hex[2:4], 16)
                        b = int(cor_hex[4:6], 16)
                        run.font.color.rgb = RGBColor(r, g, b)
            
            return True
    return False

def substituir_lista_em_placeholder(slide, placeholder_tag, lista_itens, texto_vazio="Nenhum"):
    """
    Substitui um placeholder por uma lista de itens, cada um em uma linha.
    
    Args:
        slide: Slide do python-pptx
        placeholder_tag: Nome do placeholder (sem chaves)
        lista_itens: Lista de strings ou string separada por vírgulas
        texto_vazio: Texto a ser exibido se a lista for vazia
    """
    # Converter para lista se for string
    if isinstance(lista_itens, str):
        # Remover colchetes e aspas se existirem
        if lista_itens.startswith('[') and lista_itens.endswith(']'):
            # Remover colchetes e dividir por vírgula
            lista_itens = lista_itens[1:-1]
            # Dividir por vírgula, removendo aspas extras
            itens = [item.strip().strip("'\"") for item in lista_itens.split(',') if item.strip()]
        else:
            # Se for string simples, tratar como lista de um item
            itens = [lista_itens.strip()] if lista_itens.strip() else []
    elif isinstance(lista_itens, list):
        itens = lista_itens
    else:
        itens = []
    
    # Remover itens vazios
    itens = [item for item in itens if str(item).strip()]
    
    # Se não houver itens, usar texto_vazio
    if not itens:
        texto_final = texto_vazio
        num_itens = 1
    else:
        # Juntar com quebras de linha
        texto_final = '\n'.join(itens)
        num_itens = len(itens)
    
    # Procurar pelo placeholder no slide
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
            
        if not shape.has_text_frame:
            continue
            
        # Procurar pelo placeholder no texto
        if f"{{{{{placeholder_tag}}}}}" in shape.text:
            # Limpar o texto atual
            shape.text = ""
            
            # Adicionar o novo texto
            text_frame = shape.text_frame
            text_frame.clear()  # Limpar todos os parágrafos
            
            # Configurar as propriedades do text_frame
            text_frame.word_wrap = True
            
            # Adicionar parágrafos para cada item
            if num_itens > 1:
                for i, item in enumerate(itens):
                    p = text_frame.add_paragraph()
                    p.text = str(item)
                    p.alignment = PP_ALIGN.CENTER  # Centralizar parágrafo
                    p.space_before = Pt(0)
                    
                    # Formatar o run
                    for run in p.runs:
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(0xC4, 0x80, 0x3F)  # #C4803F
                        run.font.name = 'Nunito'
                    
                    # Não adicionar espaço após o último parágrafo
                    if i < num_itens - 1:
                        p.space_after = Pt(6)  # Espaço entre linhas
            else:
                # Apenas um item (ou texto_vazio)
                p = text_frame.add_paragraph()
                p.text = texto_final
                p.alignment = PP_ALIGN.CENTER  # Centralizar parágrafo
                p.space_before = Pt(0)


                for run in p.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = RGBColor(0xC4, 0x80, 0x3F)
                    run.font.name = 'Nunito'
            
            # Configurar margens e alinhamento
            text_frame.margin_left = 0
            text_frame.margin_right = 0
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            
            # Ajustar o alinhamento vertical se necessário
            shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            break  # Sai do loop após encontrar e substituir


def calcular_media_competencias(df, COLUNAS_GERAIS):
    """
    Retorna vetor com a média das competências.
    """
    return df[COLUNAS_GERAIS].mean().values


def radar_avaliado(df_pessoa, COLUNAS_GERAIS):
    auto = df_pessoa[df_pessoa['Tipo'] == 'autoavaliacao']
    lideres = df_pessoa[df_pessoa['Tipo'] == 'avaliacao_pelo_lider']
    liderados = df_pessoa[df_pessoa['Tipo'] == 'avaliacao_pelo_liderado']

    return (
        calcular_media_competencias(auto, COLUNAS_GERAIS) +
        calcular_media_competencias(lideres, COLUNAS_GERAIS) +
        calcular_media_competencias(liderados, COLUNAS_GERAIS)
    ) / 3


def radar_frente(df_total, frente, COLUNAS_GERAIS):
    df_frente = df_total[df_total['Frente de atuação'] == frente]

    auto = df_frente[df_frente['Tipo'] == 'autoavaliacao']
    lideres = df_frente[df_frente['Tipo'] == 'avaliacao_pelo_lider']
    liderados = df_frente[df_frente['Tipo'] == 'avaliacao_pelo_liderado']

    return (
        calcular_media_competencias(auto, COLUNAS_GERAIS) +
        calcular_media_competencias(lideres, COLUNAS_GERAIS) +
        calcular_media_competencias(liderados, COLUNAS_GERAIS)
    ) / 3


def radar_geral(df_total, COLUNAS_GERAIS):
    auto = df_total[df_total['Tipo'] == 'autoavaliacao']
    lideres = df_total[df_total['Tipo'] == 'avaliacao_pelo_lider']
    liderados = df_total[df_total['Tipo'] == 'avaliacao_pelo_liderado']

    return (
        calcular_media_competencias(auto, COLUNAS_GERAIS) +
        calcular_media_competencias(lideres, COLUNAS_GERAIS) +
        calcular_media_competencias(liderados, COLUNAS_GERAIS)
    ) / 3


def calcular_nota_geral(nota_auto, nota_lider, nota_liderado, ndigits=1):
    notas = [nota_auto, nota_lider, nota_liderado]
    notas_validas = [n for n in notas if pd.notna(n)]

    if not notas_validas:
        return None

    return round(sum(notas_validas) / len(notas_validas), ndigits)


def calcular_nota_lider_geral(nota_lider_auto, nota_lider_lider, nota_lider_liderado, ndigits=1):
    notas = [nota_lider_auto, nota_lider_lider, nota_lider_liderado]
    notas_validas = [n for n in notas if pd.notna(n)]

    if not notas_validas:
        return None

    return round(sum(notas_validas) / len(notas_validas), ndigits)



# Função alternativa que mantém mais fiel à original mas calcula média das 3 fontes
def criar_radar_comparativo(df_total, nome_pessoa):
    """Versão alternativa que calcula média das 3 fontes separadamente"""
    df_pessoa = filtrar_dados_pessoa(df_total, nome_pessoa)
    
    if df_pessoa.empty:
        return None
    
    frente_pessoa = df_pessoa['Frente de atuação'].iloc[0]
    df_frente = df_total[df_total['Frente de atuação'] == frente_pessoa]
    
    categorias = {
        'Visão Sistêmica': [
            'Visao Sistemica_Identifica desafios',
            'Visao Sistematica_Compartilha ferramentas',
            'Visao Sistemica_Colabora com as demais frentes',
            'Visao Sistemica_Compreensao do impacto'
        ],
        'Gestão e Liderança': [
            'Gestao e Lideranca_Canal aberto',
            'Gestao e Lideranca_Incentiva e colabora',
            'Gestao e Lideranca_Autogestao de tempo'
        ],
        'Relacionamento': [
            'Relacionamento_Cativa parceiros',
            'Relacionamento_Cria vinculos'
        ],
        'Comunicação': [
            'Comunicacao_Assertivo',
            'Comunicacao_Nao Violenta',
            'Comunicacao_Postura Profissional'
        ],
        'Aprendizagem e Desenvolvimento': [
            'Aprendizagem_Compartilha experiencia',
            'Aprendizagem_Conhecimentos de fontes internas e externas',
            'Aprendizagem_Busca se desenvolver',
            'Aprendizagem_Autonomia'
        ],
        'Execução': [
            'Execucao_Resolve as dificuldades',
            'Execucao_Imprevistos e alteracoes',
            'Execucao_Ideias em Acoes'
        ]
    }
    
    # Função para calcular média considerando as 3 fontes separadamente
    def calcular_media_tres_fontes(df, categorias_dict):
        medias_categoria = []
        
        for categoria, questoes in categorias_dict.items():
            questoes_existentes = [q for q in questoes if q in df.columns]
            
            if questoes_existentes:
                notas_por_fonte = {'autoavaliacao': [], 'avaliacao_pelo_lider': [], 'avaliacao_pelo_liderado': []}
                
                # Coletar notas por fonte
                for tipo in notas_por_fonte.keys():
                    df_tipo = df[df['Tipo'] == tipo]
                    if not df_tipo.empty:
                        for questao in questoes_existentes:
                            notas = df_tipo[questao].dropna()
                            if not notas.empty:
                                notas_por_fonte[tipo].extend(notas.tolist())
                
                # Calcular média de cada fonte
                medias_fontes = []
                for tipo, notas in notas_por_fonte.items():
                    if notas:
                        medias_fontes.append(sum(notas) / len(notas))
                
                # Calcular média final (média das médias das fontes disponíveis)
                if medias_fontes:
                    media_final = sum(medias_fontes) / len(medias_fontes)
                    medias_categoria.append(media_final)
                else:
                    medias_categoria.append(0)
            else:
                medias_categoria.append(0)
        
        return medias_categoria
    
    # Calcular médias
    medias_pessoa = calcular_media_tres_fontes(df_pessoa, categorias)
    medias_frente = calcular_media_tres_fontes(df_frente, categorias)
    medias_total = calcular_media_tres_fontes(df_total, categorias)
    
    # Resto do código igual à primeira versão...
    # ... (configuração do gráfico, plotagem, etc.)
    
    # Configurar radar
    N = len(categorias)
    angulos = [n / float(N) * 2 * np.pi for n in range(N)]
    angulos += angulos[:1]
    
    medias_pessoa_fechado = medias_pessoa + [medias_pessoa[0]]
    medias_frente_fechado = medias_frente + [medias_frente[0]]
    medias_total_fechado = medias_total + [medias_total[0]]
    
    fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(projection='polar'))
    
    cores = {
        'Pessoa': '#1f77b4',
        'Frente': '#ff7f0e',
        'Total': '#2ca02c'
    }
    
    ax.plot(angulos, medias_pessoa_fechado, 'o-', linewidth=2, 
            label=f'{nome_pessoa}', color=cores['Pessoa'])
    ax.fill(angulos, medias_pessoa_fechado, alpha=0.1, color=cores['Pessoa'])
    
    ax.plot(angulos, medias_frente_fechado, 'o-', linewidth=2, 
            label=f'Frente: {frente_pessoa}', color=cores['Frente'])
    ax.fill(angulos, medias_frente_fechado, alpha=0.1, color=cores['Frente'])
    
    ax.plot(angulos, medias_total_fechado, 'o-', linewidth=2, 
            label='Projeto GA-RN', color=cores['Total'])
    ax.fill(angulos, medias_total_fechado, alpha=0.1, color=cores['Total'])
    
    # Configuração do gráfico...
    ax.set_xticks(angulos[:-1])
    ax.set_xticklabels(list(categorias.keys()), fontsize=11)
    
    # Ajustar labels (mesmo código da original)...
    for label, angle in zip(ax.get_xticklabels(), angulos[:-1]):
        texto = label.get_text()
        if 0 <= angle < np.pi:
            alinhamento = 'left'
        else:
            alinhamento = 'right'
        if texto == 'Relacionamento':
            alinhamento = 'right'
        if texto == 'Execução':
            alinhamento = 'left'
        label.set_horizontalalignment(alinhamento)
        label.set_rotation(angle * 180 / np.pi - 90)
    
    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(['1', '2', '3', '4', '5'], fontsize=9)
    ax.grid(True, alpha=0.3)
    
    ax.legend(loc='upper left', bbox_to_anchor=(-0.35, 1.25), 
              fontsize=10, framealpha=0.9)
    
    
    caminho_grafico = os.path.join(DIRETORIO_GRAFICOS_TEMP, f"radar_comparativo_v2_{nome_pessoa.replace(' ', '_')}.png")
    fig.savefig(caminho_grafico, dpi=150, bbox_inches='tight', pad_inches=0.5)
    plt.close(fig)
    
    return caminho_grafico


def gerar_relatorio_pessoa(nome_pessoa, df_total):
    """Gera relatório para uma pessoa específica"""
    print(f"Gerando relatório para: {nome_pessoa}")
    
    # Filtrar dados da pessoa
    df_pessoa = filtrar_dados_pessoa(df_total, nome_pessoa)
    
    if df_pessoa.empty:
        print(f"  ⚠️  Nenhum dado encontrado para {nome_pessoa}")
        return None
    
    # Carregar template
    prs = Presentation(CAMINHO_TEMPLATE)



    # =================== SLIDE 1: CAPA ===================
    slide0 = prs.slides[0]

   # Procurar pelo placeholder {{nome}} e formatá-lo
    for shape in slide0.shapes:
        if hasattr(shape, "text") and "{{nome}}" in shape.text:
            # Substituir o texto
            shape.text = shape.text.replace("{{nome}}", nome_pessoa)
            
            # Formatar a fonte - tamanho 18 e negrito
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)  # Tamanho 18
                    run.font.bold = True    # Negrito
            
            break  # Parar após encontrar o primeiro placeholder {{nome}}


# =================== SLIDE 1: AVALIAÇÕES RECEBIDAS ===================
    slide1 = prs.slides[1]  # Slide 2 (índice 1)

    # Obter lideres e liderados (qualquer linha serve, mas garantir que não seja NaN)
    linha_base = df_pessoa.iloc[0]

    # Tratar valores NaN
    lista_lideres = linha_base['Lideres'] if pd.notna(linha_base['Lideres']) else ""
    lista_liderados = linha_base['Liderados'] if pd.notna(linha_base['Liderados']) else ""

    # Função para limpar e formatar a lista
    def formatar_lista(lista_str):
        if not lista_str or pd.isna(lista_str):
            return []
        
        # Se for string que parece lista
        if isinstance(lista_str, str):
            if lista_str.startswith('[') and lista_str.endswith(']'):
                # Remover colchetes e dividir
                conteudo = lista_str[1:-1]
                # Dividir por vírgula e limpar
                itens = [item.strip().strip("'\"") for item in conteudo.split(',') if item.strip()]
                return itens
            elif ',' in lista_str:
                # Dividir por vírgula
                return [item.strip() for item in lista_str.split(',') if item.strip()]
            else:
                # Único item
                return [lista_str.strip()] if lista_str.strip() else []
        elif isinstance(lista_str, list):
            return [str(item).strip() for item in lista_str if str(item).strip()]
        return []

    # Formatar as listas
    lideres_formatado = formatar_lista(lista_lideres)
    liderados_formatado = formatar_lista(lista_liderados)

    # Substituir placeholders com a nova função
    substituir_lista_em_placeholder(
        slide1,
        "lideres",
        lideres_formatado,
        texto_vazio="Não houve avaliações de líderes"
    )

    substituir_lista_em_placeholder(
        slide1,
        "liderados",
        liderados_formatado,
        texto_vazio="Não houve avaliações de liderados"
    )



    # =================== SLIDE 3: COMPETÊNCIAS GERAIS ===================
    slide2 = prs.slides[2]
    
    cor_nota = "#CC8A42"

    # Calcular médias gerais
    nota_auto = calcular_medias_tipo(df_pessoa, 'autoavaliacao', COLUNAS_GERAIS)
    nota_lider = calcular_medias_tipo(df_pessoa, 'avaliacao_pelo_lider', COLUNAS_GERAIS)
    nota_liderado = calcular_medias_tipo(df_pessoa, 'avaliacao_pelo_liderado', COLUNAS_GERAIS)


    nota_geral = calcular_nota_geral(
        nota_auto.mean() if nota_auto is not None else None,
        nota_lider.mean() if nota_lider is not None else None,
        nota_liderado.mean() if nota_liderado is not None else None
    )

    # Substituir placeholders de notas
    substituir_texto_formatado(slide2, "nota_auto", 
                            formatar_nota(nota_auto.mean() if nota_auto is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=True)
    substituir_texto_formatado(slide2, "nota_lider", 
                            formatar_nota(nota_lider.mean() if nota_lider is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=True)
    substituir_texto_formatado(slide2, "nota_liderado", 
                            formatar_nota(nota_liderado.mean() if nota_liderado is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=True)
    substituir_texto_formatado(slide2, "nota_geral", 
                            formatar_nota(nota_geral.mean() if nota_geral is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=False)   
    
    # Gerar e adicionar radar geral
    radar_geral_path = criar_radar_geral(df_pessoa)
    if radar_geral_path:
        # Substituir placeholder de radar
        for shape in slide2.shapes:
            if hasattr(shape, "text") and "{{radar_geral}}" in shape.text:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Remover placeholder
                sp = shape._element
                sp.getparent().remove(sp)
                
                # Adicionar imagem do radar
                slide2.shapes.add_picture(radar_geral_path, left, top, width, height)
                break
    
    # =================== SLIDE 4: COMPETÊNCIAS DE LIDERANÇA ===================
    slide3 = prs.slides[3]

    # Calcular médias de liderança
    nota_lider_auto = calcular_medias_tipo(df_pessoa, 'autoavaliacao', COLUNAS_LIDERANCA)
    nota_lider_lider = calcular_medias_tipo(df_pessoa, 'avaliacao_pelo_lider', COLUNAS_LIDERANCA)
    nota_lider_liderado = calcular_medias_tipo(df_pessoa, 'avaliacao_pelo_liderado', COLUNAS_LIDERANCA)


    nota_lider_geral = calcular_nota_geral(
        nota_lider_auto.mean() if nota_lider_auto is not None else None,
        nota_lider_lider.mean() if nota_lider_lider is not None else None,
        nota_lider_liderado.mean() if nota_lider_liderado is not None else None
    )

    # Substituir placeholders
    substituir_texto_formatado(slide3, "nota_lider_auto", 
                            formatar_nota(nota_lider_auto.mean() if nota_lider_auto is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=True)
    substituir_texto_formatado(slide3, "nota_lider_lider", 
                            formatar_nota(nota_lider_lider.mean() if nota_lider_lider is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=True)
    substituir_texto_formatado(slide3, "nota_lider_liderado", 
                            formatar_nota(nota_lider_liderado.mean() if nota_lider_liderado is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=True)
    substituir_texto_formatado(slide3, "nota_lider_geral", 
                            formatar_nota(nota_lider_geral.mean() if nota_lider_geral is not None else None),
                            tamanho_fonte=14, cor_hex=cor_nota, centralizar=False)

    # Gerar e adicionar radar de liderança
    radar_lideranca_path = criar_radar_lideranca(df_pessoa)
    if radar_lideranca_path:
        # Substituir placeholder de radar (segundo slide)
        for shape in slide3.shapes:
            if hasattr(shape, "text") and "{{radar_geral}}" in shape.text:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Remover placeholder
                sp = shape._element
                sp.getparent().remove(sp)
                
                # Adicionar imagem do radar
                slide3.shapes.add_picture(radar_lideranca_path, left, top, width, height)
                break

    # =================== SLIDE 5: COMPARATIVO: COMPETÊNCIAS GERAIS ===================
    slide4 = prs.slides[4]

    # Gerar e adicionar radar comparativo
    radar_comparativo_path = criar_radar_comparativo(df_total, nome_pessoa)

    if radar_comparativo_path:
        for shape in slide4.shapes:
            if hasattr(shape, "text") and "{{radar_geral_comparativo}}" in shape.text:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Remover placeholder
                sp = shape._element
                sp.getparent().remove(sp)
                
                # Adicionar imagem do radar
                slide4.shapes.add_picture(radar_comparativo_path, left, top, width, height)
                break



    # =================== SLIDES 7-...: DETALHAMENTO DAS COMPETÊNCIAS ===================
    # Para cada questão, gerar gráfico e substituir no slide correspondente
    for questao, placeholder in MAPEAMENTO_QUESTOES_GRAFICOS.items():
        # Encontrar qual slide tem este placeholder
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and f"{{{{{placeholder}}}}}" in shape.text:
                    # Gerar gráfico
                    grafico_path = criar_grafico_barras_questao(df_pessoa, questao, placeholder)
                    
                    if grafico_path:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # Remover placeholder
                        sp = shape._element
                        sp.getparent().remove(sp)
                        
                        # Adicionar imagem do gráfico
                        slide.shapes.add_picture(grafico_path, left, top, width, height)
                    
                    break
    



 
    # =================== SLIDES FINAIS: QUESTÕES ABERTAS ===================

    def extrair_respostas_abertas_por_tipo(df_pessoa):
        """Extrai todas as respostas abertas para cada tipo, mantendo todas as linhas"""
        respostas_por_tipo = {}
        
        for tipo in ['autoavaliacao', 'avaliacao_pelo_lider', 'avaliacao_pelo_liderado']:
            df_tipo = df_pessoa[df_pessoa['Tipo'] == tipo]
            
            # Inicializar listas vazias
            pontos_fortes = []
            oportunidades = []
            
            if not df_tipo.empty:
                # Extrair TODAS as respostas (não apenas únicas)
                for _, row in df_tipo.iterrows():
                    # Pontos Fortes
                    pf = row['Pontos Fortes']
                    if pd.notna(pf) and str(pf).strip():
                        pontos_fortes.append(str(pf).strip())
                    
                    # Oportunidades de Desenvolvimento
                    od = row['Oportunidades de Desenvolvimento']
                    if pd.notna(od) and str(od).strip():
                        oportunidades.append(str(od).strip())
            
            # Se não houver respostas, usar texto padrão
            respostas_por_tipo[tipo] = {
                'PF': pontos_fortes if pontos_fortes else ["Não houve resposta"],
                'OD': oportunidades if oportunidades else ["Não houve resposta"]
            }
        
        return respostas_por_tipo

    # Extrair respostas abertas
    respostas_por_tipo = extrair_respostas_abertas_por_tipo(df_pessoa)

    # Mapeamento CORRETO dos placeholders para cada slide
    # Slide 27: PF_auto, Slide 28: PF_lider, Slide 29: PF_liderado
    # Slide 30: OD_auto, Slide 31: OD_lider, Slide 32: OD_liderado
    mapeamento_placeholders = [
        (26, "PF_auto", respostas_por_tipo['autoavaliacao']['PF']),
        (27, "PF_lider", respostas_por_tipo['avaliacao_pelo_lider']['PF']),
        (28, "PF_liderado", respostas_por_tipo['avaliacao_pelo_liderado']['PF']),
        (29, "OD_auto", respostas_por_tipo['autoavaliacao']['OD']),
        (30, "OD_lider", respostas_por_tipo['avaliacao_pelo_lider']['OD']),
        (31, "OD_liderado", respostas_por_tipo['avaliacao_pelo_liderado']['OD'])
    ]

    print(f"Total de slides na apresentação: {len(prs.slides)}")

    # Função para substituir texto com formatação específica
    def substituir_texto_formatado2(slide, placeholder_tag, lista_respostas, cor_hex="#C4803F"):
        """
        Substitui um placeholder por uma lista de respostas, cada uma em um parágrafo.
        """
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
                
            if not shape.has_text_frame:
                continue
                
            # Verificar se o placeholder está no texto da shape
            texto_shape = shape.text
            
            if f"{{{{{placeholder_tag}}}}}" in texto_shape:
                print(f"  Encontrado placeholder {{{{{placeholder_tag}}}}} no slide {prs.slides.index(slide)}")
                
                # Limpar o texto atual
                shape.text = ""
                
                # Adicionar o novo texto
                text_frame = shape.text_frame
                text_frame.clear()
                
                # Converter cor hex para RGB
                cor_rgb = RGBColor(
                    int(cor_hex[1:3], 16),
                    int(cor_hex[3:5], 16),
                    int(cor_hex[5:7], 16)
                )
                
                # Adicionar cada resposta como um parágrafo separado
                for i, resposta in enumerate(lista_respostas):
                    # Verificar se a resposta não é vazia
                    if not resposta or resposta == "":
                        continue
                        
                    # Adicionar parágrafo
                    p = text_frame.add_paragraph()
                    p.text = resposta
                    
                    # Formatar o parágrafo
                    p.alignment = PP_ALIGN.LEFT
                    p.space_before = Pt(0)
                    
                    # Formatar cada run no parágrafo
                    for run in p.runs:
                        run.font.size = Pt(9)  # Tamanho ligeiramente maior para melhor leitura
                        run.font.color.rgb = cor_rgb
                        run.font.name = 'Nunito'
                    
                    # Adicionar espaço após o parágrafo (exceto no último)
                    if i < len(lista_respostas) - 1:
                        p.space_after = Pt(8)
                
                return True
        
        return False

    # Substituir todos os placeholders
    for slide_idx, placeholder, lista_respostas in mapeamento_placeholders:
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            print(f"Processando slide {slide_idx}: placeholder {{{{{placeholder}}}}} com {len(lista_respostas)} respostas")
            
            # Verificar se há respostas
            if lista_respostas and len(lista_respostas) > 0:
                print(f"  Primeira resposta: {lista_respostas[0][:50]}..." if len(lista_respostas[0]) > 50 else f"  Primeira resposta: {lista_respostas[0]}")
            
            # Substituir o placeholder
            sucesso = substituir_texto_formatado2(slide, placeholder, lista_respostas)
            
            if sucesso:
                print(f"  ✓ Placeholder substituído com sucesso")
            else:
                print(f"  ✗ Placeholder NÃO encontrado no slide")
        else:
            print(f"AVISO: Slide {slide_idx} não existe na apresentação")

    # Função para limpar placeholders não substituídos
    def limpar_placeholders_restantes(slide):
        """Remove placeholders não substituídos"""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                texto_original = shape.text
                
                # Verificar se há placeholders no formato {{...}}
                if "{{" in texto_original and "}}" in texto_original:
                    # Usar regex para encontrar todos os placeholders
                    import re
                    # Encontrar todos os padrões {{...}}
                    placeholders = re.findall(r'\{\{[^{}]+\}\}', texto_original)
                    
                    if placeholders:
                        # Substituir cada placeholder por string vazia
                        for ph in placeholders:
                            texto_original = texto_original.replace(ph, '')
                        
                        # Atualizar o texto da shape
                        shape.text = texto_original
                        
                        # Limpar formatação se o texto ficou vazio
                        if texto_original.strip() == "":
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = ""

    # Limpar placeholders restantes em todos os slides
    for slide in prs.slides:
        limpar_placeholders_restantes(slide)

    print("Processamento das questões abertas concluído!")







    
    # Salvar apresentação
    nome_arquivo = f"Relatorio_{nome_pessoa.replace(' ', '_').replace('/', '_')}.pptx"
    caminho_saida = os.path.join(DIRETORIO_SAIDA, nome_arquivo)
    prs.save(caminho_saida)
    
    print(f"  ✅ Relatório salvo em: {caminho_saida}")
    
    # Limpar gráficos temporários
    limpar_arquivos_temporarios(df_pessoa)
    
    return caminho_saida

def limpar_arquivos_temporarios(df_pessoa):
    """Limpa arquivos temporários gerados para uma pessoa"""
    import glob
    import os
    
    # Padrão de arquivos temporários
    padrao = os.path.join(DIRETORIO_GRAFICOS_TEMP, f"*_{id(df_pessoa)}.png")
    
    for arquivo in glob.glob(padrao):
        try:
            os.remove(arquivo)
        except:
            pass



# =================== FUNÇÃO PRINCIPAL ===================
def gerar_relatorios_todos(df_total):
    """Gera relatórios para todas as pessoas no dataframe"""
    print("=" * 60)
    print("INICIANDO GERAÇÃO DE RELATÓRIOS")
    print("=" * 60)
    
    # Verificar se há dados
    if df_total.empty:
        print("Dataframe vazio!")
        return
    
    # Listar pessoas únicas
    pessoas = df_total['Nome'].unique()
    print(f"Total de pessoas encontradas: {len(pessoas)}")
    
    # Gerar relatório para cada pessoa
    relatorios_gerados = []
    
    for i, nome_pessoa in enumerate(pessoas, 1):
        print(f"\n[{i}/{len(pessoas)}] ", end="")
        caminho_relatorio = gerar_relatorio_pessoa(nome_pessoa, df_total)
        if caminho_relatorio:
            relatorios_gerados.append(caminho_relatorio)
    
    # Resumo
    print("\n" + "=" * 60)
    print("RESUMO DA GERAÇÃO")
    print("=" * 60)
    print(f"Total de pessoas processadas: {len(pessoas)}")
    print(f"Relatórios gerados com sucesso: {len(relatorios_gerados)}")
    print(f"Diretório de saída: {DIRETORIO_SAIDA}")
    
    return relatorios_gerados

# =================== EXECUÇÃO ===================
if __name__ == "__main__":
    
    # Verificar estrutura do dataframe
    print("Colunas disponíveis no dataframe:")
    print(df_total.columns.tolist())
    print(f"\nTotal de registros: {len(df_total)}")
    print(f"Tipos de avaliação: {df_total['Tipo'].unique()}")
    
    # Gerar todos os relatórios
    relatorios = gerar_relatorios_todos(df_total)
    
    # Limpar diretório de gráficos temporários completamente
    try:
        import shutil
        shutil.rmtree(DIRETORIO_GRAFICOS_TEMP)
        os.makedirs(DIRETORIO_GRAFICOS_TEMP, exist_ok=True)
    except:
        pass




